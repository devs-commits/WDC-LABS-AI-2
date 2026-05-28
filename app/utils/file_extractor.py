"""
File extraction utility - handles multiple file formats
Upgraded for large dataset safety and Excel structure preservation.
"""

import io
import json
import csv
import PyPDF2
import urllib.parse
from docx import Document

try:
    import openpyxl
    EXCEL_SUPPORT = True
except ImportError:
    EXCEL_SUPPORT = False

try:
    from pptx import Presentation
    PPTX_SUPPORT = True
except ImportError:
    PPTX_SUPPORT = False


def extract_text_from_file(file_url: str, file_content_bytes: bytes, mime_type: str = None) -> str:
    """
    Extract text from various file formats.
    Supports: PDF, DOCX, XLSX, PPTX, CSV, TXT, JSON, Images
    """
    if not file_content_bytes:
        return "[Error: Empty file received from server]"

    try:
        # 1. Clean the URL to get the TRUE extension (removes ?tokens=123)
        clean_url = urllib.parse.urlparse(file_url).path if file_url else ""
        file_ext = clean_url.lower().split('.')[-1] if '.' in clean_url else ""
        
        # 2. Make mime_type safe
        mime = (mime_type or "").lower()

        # --- ROUTING LOGIC ---

        # PDF files
        if file_ext == "pdf" or "pdf" in mime:
            return extract_pdf(file_content_bytes)
        
        # Excel files (.xlsx, .xls)
        elif file_ext in ["xlsx", "xls"] or "sheet" in mime or "excel" in mime:
            result = extract_excel(file_content_bytes)
            # Failsafe: If it failed because it's secretly a CSV disguised as an Excel file
            if "extraction error" in result:
                csv_fallback = extract_csv(file_content_bytes)
                if "error" not in csv_fallback:
                    return csv_fallback
            return result
        
        # Word documents (.docx)
        elif file_ext in ["docx", "doc"] or "word" in mime or "document" in mime:
            return extract_docx(file_content_bytes)
        
        # PowerPoint files (.pptx)
        elif file_ext == "pptx" or "presentation" in mime:
            return extract_pptx(file_content_bytes)
        
        # CSV files
        elif file_ext == "csv" or "csv" in mime:
            return extract_csv(file_content_bytes)
        
        # JSON files
        elif file_ext == "json" or "json" in mime:
            return extract_json(file_content_bytes)
        
        # Plain text and Code files
        elif file_ext in ["txt", "md", "py", "js", "html", "css", "sql"] or "text" in mime:
            return process_text_file(file_content_bytes, file_ext)
        
        # --- MAGIC BYTE FALLBACKS FOR UNRECOGNIZED EXTENSIONS ---
        
        # 'PK' is the magic byte signature for all OpenXML (xlsx, docx) files
        elif file_content_bytes.startswith(b'PK\x03\x04'): 
            # Try Excel first (most common for Data Analytics track)
            excel_text = extract_excel(file_content_bytes)
            if "Excel extraction error" not in excel_text and "support not installed" not in excel_text:
                return excel_text
            
            # Try Word second
            docx_text = extract_docx(file_content_bytes)
            if "DOCX extraction error" not in docx_text:
                return docx_text
                
            return "[Error: Unrecognized OpenXML format. Could not parse as Excel or Word.]"
        
        # FALLBACK 2: Absolute Default try UTF-8 decode
        else:
            return process_text_file(file_content_bytes, file_ext)
    
    except Exception as e:
        return f"[Error extracting file: {str(e)}]"


def process_text_file(file_bytes: bytes, file_ext: str) -> str:
    """Helper to process plain text and code files"""
    try:
        raw_text = file_bytes.decode("utf-8", errors="ignore")
        # Inject line numbers for code files so agents can easily point out errors
        if file_ext in ["py", "js", "html", "css", "sql"]:
            lines = raw_text.split('\n')
            numbered_lines = [f"[Line {i+1}] {line}" for i, line in enumerate(lines)]
            return "\n".join(numbered_lines)
        return raw_text
    except Exception:
        return "[Binary file - cannot extract text]"


def extract_pdf(file_bytes: bytes) -> str:
    """Extract full text from PDF"""
    try:
        reader = PyPDF2.PdfReader(io.BytesIO(file_bytes))
        text = ""
        for page_num, page in enumerate(reader.pages): 
            text += f"\n--- Page {page_num + 1} ---\n"
            text += page.extract_text() or "[No text found on page]"
        return text 
    except Exception as e:
        return f"[PDF extraction error: {str(e)}]"


def extract_docx(file_bytes: bytes) -> str:
    """Extract full text from Word document (.docx)"""
    try:
        doc = Document(io.BytesIO(file_bytes))
        text = ""
        for para in doc.paragraphs:
            if para.text.strip():
                text += para.text + "\n"
        for table in doc.tables:
            text += "\n[TABLE]\n"
            for row in table.rows:
                for cell in row.cells:
                    text += cell.text + " | "
                text += "\n"
        return text
    except Exception as e:
        return f"[DOCX extraction error: {str(e)}]"


def extract_excel(file_bytes: bytes) -> str:
    """Extract text from Excel file with strict limits to prevent AI hallucination"""
    if not EXCEL_SUPPORT:
        return "[Excel extraction error: openpyxl library is not installed on the server.]"
    
    try:
        # data_only=True ensures we get calculated values, not the raw formula strings
        workbook = openpyxl.load_workbook(io.BytesIO(file_bytes), data_only=True)
        text = ""
        
        for sheet_name in workbook.sheetnames: 
            worksheet = workbook[sheet_name]
            text += f"\n--- Sheet: {sheet_name} ---\n"
            
            row_count = 0
            for row in worksheet.iter_rows(values_only=True): 
                # Convert to string and handle None types. Strip newlines inside cells to prevent layout breaking.
                row_text = " | ".join(str(cell).replace('\n', ' ').replace('\r', '').strip() if cell is not None else "" for cell in row)
                
                # Only add if row is not completely empty
                if row_text.replace(" | ", "").strip(): 
                    text += row_text + "\n"
                    row_count += 1
                
                # Hard limit to 1500 rows to prevent blowing out the LLM's context memory!
                if row_count >= 1500:
                    text += "\n[... TRUNCATED: Exceeded 1500 rows to prevent system memory overload ...]\n"
                    break
                    
        if not text.strip():
            return "[Excel file appears to be completely empty]"
            
        return text
        
    except Exception as e:
        error_msg = str(e).lower()
        if "does not support the old .xls" in error_msg or "invalid file" in error_msg:
            return "[Excel extraction error: This appears to be an older .xls file or an unsupported format. Please convert it to .xlsx or .csv and try again.]"
        return f"[Excel extraction error: {str(e)}]"


def extract_pptx(file_bytes: bytes) -> str:
    """Extract full text from PowerPoint presentation (.pptx)"""
    if not PPTX_SUPPORT:
        return "[PowerPoint support not installed - install python-pptx]"
    try:
        prs = Presentation(io.BytesIO(file_bytes))
        text = ""
        for slide_num, slide in enumerate(prs.slides): 
            text += f"\n--- Slide {slide_num + 1} ---\n"
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    text += shape.text + "\n"
        return text
    except Exception as e:
        return f"[PPTX extraction error: {str(e)}]"


def extract_csv(file_bytes: bytes) -> str:
    """Extract full text from CSV file"""
    try:
        text_content = file_bytes.decode("utf-8", errors="ignore")
        csv_reader = csv.reader(io.StringIO(text_content))
        text = ""
        row_count = 0
        
        for row in csv_reader:
            text += " | ".join(row) + "\n"
            row_count += 1
            if row_count >= 1500:
                text += "\n[... TRUNCATED: Exceeded 1500 rows to prevent system memory overload ...]\n"
                break
                
        return text
    except Exception as e:
        return f"[CSV extraction error: {str(e)}]"


def extract_json(file_bytes: bytes) -> str:
    """Extract full text from JSON file"""
    try:
        json_content = json.loads(file_bytes.decode("utf-8"))
        return json.dumps(json_content, indent=2)
    except Exception as e:
        return f"[JSON extraction error: {str(e)}]"